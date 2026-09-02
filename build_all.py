import os
import shutil
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import docx
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Universally safe high-contrast colors (100% identical on Mobile Office & Laptop)
    BG_COLOR = RGBColor(10, 19, 14)       # #0A130E Solid Forest Dark
    CARD_BG = RGBColor(20, 36, 26)        # #14241A Solid Card Surface
    ACCENT_GREEN = RGBColor(34, 197, 94)  # #22C55E Bright Neon Emerald
    ACCENT_GOLD = RGBColor(245, 158, 11)  # #F59E0B Solar Gold
    ACCENT_CYAN = RGBColor(6, 182, 212)   # #06B6D4 Electric Cyan
    ACCENT_RED = RGBColor(239, 68, 68)    # #EF4444 Alert Coral
    TEXT_LIGHT = RGBColor(255, 255, 255)  # #FFFFFF Pure Crisp White
    TEXT_MUTED = RGBColor(226, 232, 240)  # #E2E8F0 Bright Silver-Slate
    TEXT_HIGHLIGHT = RGBColor(134, 239, 172) # #86EFAC Soft Light Green

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        # 1. Native slide background for mobile PowerPoint / Google Slides app
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # 2. Solid backdrop rectangle
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()

    def add_card(slide, left, top, width, height, border_color=ACCENT_GREEN, bg_color=CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.75)
        return card

    def add_header(slide, title_text, category_text="WEED MANAGEMENT", slide_num=None):
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = ACCENT_GREEN
        top_line.line.fill.background()

        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.0), Inches(0.35))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = f"• {category_text.upper()}"
        p_tag.font.name = "Calibri"
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_GREEN

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10.5), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Calibri"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(11.3), Inches(0.35), Inches(1.5), Inches(0.4))
            tf_num = num_box.text_frame
            p_num = tf_num.paragraphs[0]
            p_num.text = f"Slide {slide_num:02d} / 12"
            p_num.alignment = PP_ALIGN.RIGHT
            p_num.font.name = "Calibri"
            p_num.font.size = Pt(11)
            p_num.font.bold = True
            p_num.font.color.rgb = ACCENT_GOLD

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    top_line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = ACCENT_GREEN
    top_line.line.fill.background()

    add_card(s1, Inches(0.8), Inches(1.2), Inches(6.8), Inches(5.5), border_color=ACCENT_GREEN)
    tbox = s1.shapes.add_textbox(Inches(1.1), Inches(1.5), Inches(6.2), Inches(5.0))
    tf = tbox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "AGRICULTURE & CROP SCIENCE"
    p0.font.name = "Calibri"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "Weed Management\nin Agriculture"
    p1.font.name = "Calibri"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_LIGHT
    p1.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Key Topics Covered:"
    p2.font.name = "Calibri"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.space_after = Pt(6)

    topics = [
        "Meaning of Weed & Weed Management",
        "Importance & Core Principles",
        "Preventive, Cultural, Mechanical & Biological Methods",
        "Mulching Methods (Organic & Plastic)",
        "Advantages & Practical Limitations"
    ]
    for top in topics:
        pt = tf.add_paragraph()
        pt.text = f"▪  {top}"
        pt.font.name = "Calibri"
        pt.font.size = Pt(12)
        pt.font.color.rgb = TEXT_MUTED
        pt.space_after = Pt(4)

    img1 = "assets/integrated_weed_management.jpg"
    if os.path.exists(img1):
        s1.shapes.add_picture(img1, Inches(7.9), Inches(1.2), Inches(4.6), Inches(5.5))

    s1.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 1 - Welcome):\n"
        "Welcome everyone. Today we are presenting a complete guide to Weed Management in Agriculture.\n"
        "We will go through the definition of weeds, why managing them is critical for crop production, core principles, all key control methods—including preventive, cultural, mechanical, mulching, and biological—and discuss their advantages and limitations."
    )

    # =========================================================================
    # SLIDE 2: Meaning of Weeds & Weed Management
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "Meaning: What is a Weed & Weed Management?", slide_num=2)

    img2 = "assets/weed_meaning_intro.jpg"
    if os.path.exists(img2):
        s2.shapes.add_picture(img2, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    meanings = [
        ("What is a Weed?", "A weed is any plant growing where it is not wanted. For example, wild plants that grow in crop fields and compete with crops for food and water.", ACCENT_GREEN, "🌿"),
        ("What is Weed Management?", "The practice of controlling and suppressing unwanted weeds so crops can grow healthy, receive full nutrients, and produce good harvest yields.", ACCENT_GOLD, "🎯"),
        ("How Weeds Harm Crops", "Weeds compete directly with main crops for sunlight, soil moisture, space, and fertilizers, reducing overall harvest quality and quantity.", ACCENT_CYAN, "⚠️"),
        ("Key Farming Goal", "Keep fields clean and free from heavy weed growth, especially during early crop development.", TEXT_LIGHT, "💡")
    ]

    top_pos = 1.5
    for title, desc, color, icon in meanings:
        add_card(s2, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s2.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s2.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 2 - Meaning of Weed & Weed Management):\n"
        "Let us start with the basic definition:\n"
        "1. What is a weed? It is simply a plant growing where it is not wanted. Even a tomato plant in a wheat field is a weed because it competes with the wheat.\n"
        "2. What is weed management? It is the organized method of controlling weeds so our main crops get the water, fertilizer, and sunlight they need to produce a high yield."
    )

    # =========================================================================
    # SLIDE 3: Importance of Weed Management
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "Importance: Why Weed Management is Important", slide_num=3)

    img3 = "assets/weed_crop_competition.jpg"
    if os.path.exists(img3):
        s3.shapes.add_picture(img3, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    importance_points = [
        ("Protects Crop Yield & Harvest", "Weeds rob crops of sunlight, space, and nutrients. Controlling weeds prevents large yield losses.", ACCENT_GOLD, "📈"),
        ("Saves Water & Expensive Fertilizers", "Weeds absorb water and fertilizers very rapidly. Managing weeds ensures inputs nourish the crop.", ACCENT_GREEN, "💧"),
        ("Prevents Pests & Crop Diseases", "Weeds provide hiding and breeding grounds for harmful insects, viruses, and plant pathogens.", ACCENT_CYAN, "🛡️"),
        ("Improves Produce Quality & Income", "Clean fields yield higher quality produce free from weed seeds, resulting in better market prices.", TEXT_LIGHT, "💰")
    ]

    top_pos = 1.5
    for title, desc, color, icon in importance_points:
        add_card(s3, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s3.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s3.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 3 - Importance of Weed Management):\n"
        "Why is weed management so important for farmers?\n"
        "1. Yield protection: Without weed control, crop yield can drop heavily.\n"
        "2. Resource efficiency: Weeds drink the irrigation water and eat the costly fertilizer intended for your crops.\n"
        "3. Disease and pest prevention: Weeds host pests and fungi that attack your crops.\n"
        "4. Higher market prices: Clean grain and produce sell at higher prices in the market."
    )

    # =========================================================================
    # SLIDE 4: Core Principles of Weed Management
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Principles of Weed Management", slide_num=4)

    principles = [
        ("1. Prevention First", "Prevent weed seeds from entering the field in the first place. Keeping seeds out is much easier than removing grown weeds.", ["Use certified clean seeds.", "Clean equipment before moving to fresh fields."], ACCENT_GREEN, "🛡️"),
        ("2. Critical Early Window (Timing)", "Crops are most sensitive during early growth (first 20 to 35 days). Keeping fields weed-free during this window ensures healthy crop establishment.", ["Focus weeding efforts early.", "Crop leaves will naturally shade out late weeds."], ACCENT_GOLD, "⏱️"),
        ("3. Combine Multiple Methods", "Do not rely on just one single technique. Using cultural, mechanical, and biological tools together gives the best results.", ["Combines mechanical and natural tools.", "Prevents weeds from adapting."], ACCENT_CYAN, "🔄"),
        ("4. Stop Weed Seed Production", "Never let weeds flower and produce seeds in your field. Removing weeds before seed formation steadily reduces weeds in future seasons.", ["'One year's seeds make years of weeds.'", "Remove weeds before they flower."], TEXT_LIGHT, "📉")
    ]

    coords_4box = [
        (0.8, 1.5, 5.6, 2.5),
        (6.8, 1.5, 5.7, 2.5),
        (0.8, 4.3, 5.6, 2.5),
        (6.8, 4.3, 5.7, 2.5)
    ]

    for (title, desc, bullets, color, icon), (l, t, w, h) in zip(principles, coords_4box):
        add_card(s4, Inches(l), Inches(t), Inches(w), Inches(h), border_color=color)
        tb = s4.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.12), Inches(w - 0.3), Inches(h - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p_sub = tf.add_paragraph()
        p_sub.text = desc
        p_sub.font.name = "Calibri"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_LIGHT
        p_sub.space_after = Pt(6)

        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = f"▪ {b}"
            pb.font.name = "Calibri"
            pb.font.size = Pt(10)
            pb.font.color.rgb = TEXT_MUTED

    s4.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 4 - Principles of Weed Management):\n"
        "Here are the four key principles of weed management:\n"
        "1. Prevention: Stop weed seeds before they enter the farm.\n"
        "2. Timing: Protect the crop during its critical first 30 days.\n"
        "3. Integration: Combine different methods for better results.\n"
        "4. Seed control: Remove weeds before they flower and produce seeds."
    )

    # =========================================================================
    # SLIDE 5: Preventive Methods
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Preventive Methods: Stopping Weeds Before They Enter", slide_num=5)

    prevent_items = [
        ("Clean Certified Seeds", "Always use certified, pure seeds that are free from any weed seed contamination.", "🌾"),
        ("Clean Farm Machinery", "Wash tractors, cultivators, and harvesting equipment before moving from field to field.", "🚜"),
        ("Well-Rotted Compost & Manure", "Only apply fully decomposed manure so that any weed seeds inside have been destroyed by natural heat.", "🧪"),
        ("Clean Irrigation & Field Borders", "Keep irrigation canals and field borders clear of weeds so seeds do not wash or blow into crops.", "🌊")
    ]

    coords_col4 = [
        (0.8, 1.5, 2.7, 5.3),
        (3.8, 1.5, 2.7, 5.3),
        (6.8, 1.5, 2.7, 5.3),
        (9.8, 1.5, 2.7, 5.3)
    ]
    colors_col = [ACCENT_GREEN, ACCENT_GOLD, ACCENT_CYAN, TEXT_LIGHT]

    for (title, desc, icon), (l, t, w, h), col in zip(prevent_items, coords_col4, colors_col):
        add_card(s5, Inches(l), Inches(t), Inches(w), Inches(h), border_color=col)
        tb = s5.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.2), Inches(w - 0.3), Inches(h - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}\n{title}"
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(14)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_after = Pt(16)

        p3 = tf.add_paragraph()
        p3.text = "Benefit:\nStops weeds before they ever start growing."
        p3.font.name = "Calibri"
        p3.font.size = Pt(10.5)
        p3.font.bold = True
        p3.font.color.rgb = TEXT_HIGHLIGHT

    s5.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 5 - Preventive Methods):\n"
        "Preventive methods are the most cost-effective way to manage weeds:\n"
        "1. Certified seeds: Make sure no foreign weed seeds are planted with your crop.\n"
        "2. Clean machines: Harvesters and tractors carry seeds from field to field if not washed.\n"
        "3. Well-rotted manure: Fresh manure contains active weed seeds.\n"
        "4. Clean canals: Stops weed seeds from traveling in irrigation water."
    )

    # =========================================================================
    # SLIDE 6: Cultural Methods
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Cultural Methods: Agronomic Farming Practices", slide_num=6)

    img6 = "assets/cover_crop_mulching.jpg"
    if os.path.exists(img6):
        s6.shapes.add_picture(img6, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    cultural_items = [
        ("Stale Seedbed Preparation", "Water the soil before sowing, let weeds sprout, kill them with shallow tillage, and then sow your crop into a clean bed.", ACCENT_GREEN, "🌱"),
        ("Crop Rotation", "Alternate different crops each season to disrupt the life cycles of specific weeds.", ACCENT_GOLD, "🔄"),
        ("Proper Plant Spacing & Density", "Plant crops at the right spacing so crop leaves grow together quickly and shade out weeds.", ACCENT_CYAN, "📐"),
        ("Cover Crops & Intercropping", "Grow fast-growing cover crops between crop rows to blanket the soil and smother weeds.", TEXT_LIGHT, "🌿")
    ]

    top_pos = 1.5
    for title, desc, color, icon in cultural_items:
        add_card(s6, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s6.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s6.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 6 - Cultural Methods):\n"
        "Cultural methods use smart farming habits to suppress weeds:\n"
        "1. Stale Seedbed: Irrigating early sprouts weeds so they can be removed before the real crop is sown.\n"
        "2. Crop Rotation: Rotating crops prevents any single weed type from dominating.\n"
        "3. Spacing: Proper spacing makes crops form a dense canopy that shades out weeds.\n"
        "4. Cover Crops: Growing ground-cover plants smothers weed growth naturally."
    )

    # =========================================================================
    # SLIDE 7: Mechanical Methods
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Mechanical Methods: Physical Tools & Machinery", slide_num=7)

    img7 = "assets/mechanical_weeding_tractor.jpg"
    if os.path.exists(img7):
        s7.shapes.add_picture(img7, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    mech_items = [
        ("Hand Weeding & Hoeing", "Physically pulling or cutting weeds using simple hand tools like a hoe or khurpi around crop plants.", ACCENT_GREEN, "🛠️"),
        ("Tractor Inter-Row Cultivators", "Using tractor-mounted rotary weeders or tines to cultivate and uproot weeds between crop rows quickly.", ACCENT_GOLD, "🚜"),
        ("Mowing & Slashing", "Cutting tall weeds before they flower, especially along field borders, orchards, and plantation rows.", ACCENT_CYAN, "✂️"),
        ("Soil Solarization", "Covering moist soil with clear plastic sheets during hot sunny weeks to heat the soil and kill weed seeds.", TEXT_LIGHT, "☀️")
    ]

    top_pos = 1.5
    for title, desc, color, icon in mech_items:
        add_card(s7, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s7.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s7.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 7 - Mechanical Methods):\n"
        "Mechanical methods use physical tools to destroy weeds:\n"
        "1. Hand Weeding: Great for small farms and close to delicate plants.\n"
        "2. Tractor Cultivation: Fast and efficient for large acreage.\n"
        "3. Mowing: Stops weeds from producing flowers and seeds.\n"
        "4. Soil Solarization: Uses summer sun heat under plastic sheets to sterilize weed seeds."
    )

    # =========================================================================
    # SLIDE 8: Mulching Methods (Organic & Plastic)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Mulching Methods: Organic & Plastic Soil Covers", slide_num=8)

    img8 = "assets/organic_plastic_mulching.jpg"
    if os.path.exists(img8):
        s8.shapes.add_picture(img8, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    mulch_items = [
        ("What is Mulching?", "Mulching means covering the bare soil around crops. It blocks sunlight so weed seeds cannot sprout, while holding moisture in the soil.", ACCENT_GREEN, "🛡️"),
        ("Organic Mulch (Straw & Leaves)", "Using dry straw, dried leaves, or crop residue. Naturally decomposes over time and adds organic matter to the soil.", ACCENT_GOLD, "🌾"),
        ("Plastic Mulch (Black Film)", "Laying black or silver polyethylene sheets over raised beds with small holes for crop seedlings. Widely used for vegetables and fruits.", ACCENT_CYAN, "🖤"),
        ("Key Benefits of Mulching", "• Blocks 90%+ of weed growth.\n• Saves up to 50% of irrigation water.\n• Keeps soil cool and prevents soil erosion.", TEXT_LIGHT, "⭐")
    ]

    top_pos = 1.5
    for title, desc, color, icon in mulch_items:
        add_card(s8, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s8.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s8.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 8 - Mulching Methods):\n"
        "Mulching is one of the best non-chemical weed control methods:\n"
        "1. How it works: By covering the soil, it blocks sunlight so weed seeds cannot germinate.\n"
        "2. Organic Mulch: Made of straw or dry leaves that eventually nourish the soil.\n"
        "3. Plastic Mulch: Black plastic film commonly used in vegetable and fruit farming.\n"
        "4. Benefits: Stops weeds, saves water, and keeps the soil healthy."
    )

    # =========================================================================
    # SLIDE 9: Biological Methods
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_header(s9, "Biological Methods: Natural Weed Control", slide_num=9)

    img9 = "assets/biological_control_insect.jpg"
    if os.path.exists(img9):
        s9.shapes.add_picture(img9, Inches(0.8), Inches(1.5), Inches(4.8), Inches(5.3))

    bio_items = [
        ("Beneficial Insects", "Releasing specific insects that feed exclusively on harmful weeds without eating or damaging farm crops.", ACCENT_GREEN, "🐞"),
        ("Natural Bio-Agents & Fungi", "Using specific natural fungi or microorganisms that target and weaken particular weed species.", ACCENT_GOLD, "🍄"),
        ("Grazing Animals & Ducks", "Using ducks in rice fields or grazing animals in orchards to eat weed shoots naturally.", ACCENT_CYAN, "🦆"),
        ("Major Benefit", "100% eco-friendly, zero chemical residues, and completely safe for soil, pollinators, and people.", TEXT_LIGHT, "🌿")
    ]

    top_pos = 1.5
    for title, desc, color, icon in bio_items:
        add_card(s9, Inches(5.9), Inches(top_pos), Inches(6.6), Inches(1.2), border_color=color)
        tb = s9.shapes.add_textbox(Inches(6.05), Inches(top_pos + 0.08), Inches(6.3), Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = "Calibri"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

        top_pos += 1.35

    s9.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 9 - Biological Methods):\n"
        "Biological methods use living organisms to control weeds:\n"
        "1. Specific insects that eat only the problem weed without touching crops.\n"
        "2. Natural bio-agents and fungi that suppress target weeds.\n"
        "3. Grazing animals like ducks in rice fields to consume weed sprouts.\n"
        "4. This method is completely natural with zero chemical residues."
    )

    # =========================================================================
    # SLIDE 10: Advantages & Limitations
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_header(s10, "Comparison: Advantages & Practical Limitations", slide_num=10)

    # Left: Advantages
    add_card(s10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), border_color=ACCENT_GREEN)
    tb_adv = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.9))
    tf_adv = tb_adv.text_frame
    tf_adv.word_wrap = True

    p = tf_adv.paragraphs[0]
    p.text = "✅ Advantages of Weed Management"
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    adv_list = [
        ("Higher Crop Yields", "Increases harvest yield by giving crops full access to nutrients and light."),
        ("Saves Water & Fertilizers", "Ensures irrigation and fertilizer nourish crops, not weeds."),
        ("Better Produce Quality", "Delivers clean, weed-free harvest that commands higher market prices."),
        ("Fewer Pests & Diseases", "Removes the shelter and breeding places of harmful crop pests."),
        ("Long-Term Soil Health", "Reduces weed seed buildup in the soil for easier future farming.")
    ]

    for title, desc in adv_list:
        pt = tf_adv.add_paragraph()
        pt.text = f"✔ {title}"
        pt.font.name = "Calibri"
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_LIGHT
        
        pd = tf_adv.add_paragraph()
        pd.text = f"   ↳ {desc}"
        pd.font.name = "Calibri"
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_after = Pt(6)

    # Right: Limitations
    add_card(s10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), border_color=ACCENT_RED)
    tb_lim = s10.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    tf_lim = tb_lim.text_frame
    tf_lim.word_wrap = True

    p = tf_lim.paragraphs[0]
    p.text = "⚠️ Limitations & Practical Challenges"
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(14)

    lim_list = [
        ("High Labor Demand", "Hand weeding requires significant time and manual labor during busy seasons."),
        ("Equipment & Material Costs", "Tractor weeders and plastic mulches require upfront purchase costs."),
        ("Weather Dependence", "Heavy rains and wet soil can delay mechanical weeding operations."),
        ("Plastic Waste Management", "Plastic mulch sheets must be collected after harvest to avoid farm waste.")
    ]

    for title, desc in lim_list:
        pt = tf_lim.add_paragraph()
        pt.text = f"✖ {title}"
        pt.font.name = "Calibri"
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_GOLD
        
        pd = tf_lim.add_paragraph()
        pd.text = f"   ↳ {desc}"
        pd.font.name = "Calibri"
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_MUTED
        pd.space_after = Pt(8)

    s10.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 10 - Advantages & Limitations):\n"
        "Here is a balanced comparison:\n"
        "Advantages: Higher yields, efficient fertilizer and water use, fewer pests, and cleaner harvest.\n"
        "Limitations: Hand labor is hard to find, equipment costs money, wet weather can delay weeding, and plastic mulch needs disposal.\n"
        "Combining different methods helps overcome these limitations."
    )

    # =========================================================================
    # SLIDE 11: Conclusion
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_header(s11, "Conclusion: 5 Key Golden Takeaways", slide_num=11)

    add_card(s11, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3), border_color=ACCENT_GREEN)
    tb_c = s11.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(4.9))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True

    p = tf_c.paragraphs[0]
    p.text = "🎯 Summary for Successful Weed Management"
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(14)

    rules_summary = [
        ("1. Start Clean", "Always use certified clean seeds and sanitized equipment before planting."),
        ("2. Act Early", "Keep crops weed-free during the critical first 30 days of growth."),
        ("3. Protect the Soil", "Use organic or plastic mulch to conserve water and block weed emergence."),
        ("4. Combine Methods", "Integrate cultural, mechanical, and biological practices for best results."),
        ("5. Prevent Seed Setting", "Remove weeds before they flower to stop weed buildup in future years.")
    ]

    for title, desc in rules_summary:
        pt = tf_c.add_paragraph()
        pt.text = f"✔  {title}:"
        pt.font.name = "Calibri"
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = ACCENT_GOLD
        
        pd = tf_c.add_paragraph()
        pd.text = f"    {desc}"
        pd.font.name = "Calibri"
        pd.font.size = Pt(11.5)
        pd.font.color.rgb = TEXT_LIGHT
        pd.space_after = Pt(8)

    s11.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 11 - Conclusion):\n"
        "In conclusion, remember these five simple rules:\n"
        "1. Start clean with certified seed.\n"
        "2. Act early during the first 30 days.\n"
        "3. Protect the soil with mulch.\n"
        "4. Combine multiple methods.\n"
        "5. Never let weeds produce seeds."
    )

    # =========================================================================
    # SLIDE 12: Thank You Slide
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)

    top_line = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = ACCENT_GREEN
    top_line.line.fill.background()

    add_card(s12, Inches(0.8), Inches(1.2), Inches(6.0), Inches(5.5), border_color=ACCENT_GOLD)
    tb_ty = s12.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.4), Inches(4.5))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True

    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "🌾 THANK YOU!"
    p_ty.font.name = "Calibri"
    p_ty.font.size = Pt(36)
    p_ty.font.bold = True
    p_ty.font.color.rgb = ACCENT_GOLD
    p_ty.space_after = Pt(16)

    p_sub1 = tf_ty.add_paragraph()
    p_sub1.text = "Thank you for your time and attention."
    p_sub1.font.name = "Calibri"
    p_sub1.font.size = Pt(16)
    p_sub1.font.color.rgb = TEXT_LIGHT
    p_sub1.space_after = Pt(18)

    p_sub2 = tf_ty.add_paragraph()
    p_sub2.text = "Healthy Crops  •  Thriving Farms  •  Higher Yields"
    p_sub2.font.name = "Calibri"
    p_sub2.font.size = Pt(14)
    p_sub2.font.bold = True
    p_sub2.font.color.rgb = ACCENT_GREEN
    p_sub2.space_after = Pt(22)

    p_sub3 = tf_ty.add_paragraph()
    p_sub3.text = "💬 Questions & Discussion Welcome!"
    p_sub3.font.name = "Calibri"
    p_sub3.font.size = Pt(14)
    p_sub3.font.color.rgb = TEXT_MUTED

    img12 = "assets/conclusion_thankyou_farm.jpg"
    if os.path.exists(img12):
        s12.shapes.add_picture(img12, Inches(7.1), Inches(1.2), Inches(5.4), Inches(5.5))

    s12.notes_slide.notes_text_frame.text = (
        "PRESENTER NOTES (Slide 12 - Closing):\n"
        "Thank you very much for your attention!\n"
        "We are now open for any questions, thoughts, or discussion on weed management."
    )

    prs.save("Weed_Management_Mastery.pptx")
    print("Successfully generated Weed_Management_Mastery.pptx")


def create_docx():
    doc = docx.Document()

    for section in doc.sections:
        section.top_margin = DocxInches(1.0)
        section.bottom_margin = DocxInches(1.0)
        section.left_margin = DocxInches(1.0)
        section.right_margin = DocxInches(1.0)

    p_title = doc.add_paragraph()
    r_title = p_title.add_run("Weed Management in Agriculture")
    r_title.bold = True
    r_title.font.size = DocxPt(24)
    r_title.font.color.rgb = DocxRGBColor(21, 128, 61)

    p_sub = doc.add_paragraph()
    r_sub = p_sub.add_run("Simple, Practical Guide & Presenter Handbook")
    r_sub.font.size = DocxPt(13)
    r_sub.font.color.rgb = DocxRGBColor(100, 116, 139)

    doc.add_paragraph("―" * 45)

    sections_content = [
        ("1. Meaning of Weed & Weed Management", [
            ("What is a Weed?", "A weed is any plant growing where it is not wanted. For example, wild plants that grow in crop fields and compete with crops for food, water, and sunlight."),
            ("What is Weed Management?", "The practice of controlling unwanted weeds so crops can grow healthy, receive full soil nutrients, and produce good harvest yields.")
        ]),
        ("2. Importance of Weed Management", [
            ("Protects Crop Yield", "Weeds compete with crops for light, space, and nutrients. Controlling weeds prevents severe yield losses."),
            ("Saves Water & Fertilizer", "Weeds absorb water and fertilizer rapidly. Managing weeds ensures these expensive inputs nourish your crops."),
            ("Prevents Pests & Diseases", "Weeds serve as hiding and breeding spots for harmful insects, viruses, and fungi."),
            ("Improves Produce Quality", "Clean, weed-free harvest brings higher quality produce that sells at top market prices.")
        ]),
        ("3. Core Principles of Weed Management", [
            ("1. Prevention First", "Preventing weed seeds from entering the field is much easier and cheaper than removing grown weeds."),
            ("2. Critical Early Window", "Crops are most sensitive during the first 20 to 35 days. Keeping fields clean early allows crops to form a canopy that naturally shades out weeds."),
            ("3. Combine Multiple Methods", "Use cultural, mechanical, and biological methods together instead of relying on only one tool."),
            ("4. Stop Weed Seed Production", "Never let weeds flower and produce seeds. Removing weeds early reduces weed problems in future years.")
        ]),
        ("4. Preventive Methods", [
            ("Clean Certified Seeds", "Always sow certified crop seeds free from weed seeds."),
            ("Clean Farm Machinery", "Wash tractors and harvesting machinery before moving to clean fields."),
            ("Well-Rotted Compost", "Only use fully decomposed manure where weed seeds have been eliminated by natural heat."),
            ("Clean Field Borders", "Keep irrigation canals and borders clear so seeds do not blow or wash into fields.")
        ]),
        ("5. Cultural Methods", [
            ("Stale Seedbed", "Irrigate the soil before sowing, let weeds sprout, kill them with light raking, and then sow your crop in a clean bed."),
            ("Crop Rotation", "Rotate different crops each season to disrupt the life cycles of weeds."),
            ("Proper Plant Spacing", "Plant crops at the recommended density so the crop canopy closes quickly and shades out weeds."),
            ("Cover Crops", "Grow fast-growing cover crops between rows to blanket the soil and naturally smother weeds.")
        ]),
        ("6. Mechanical Methods", [
            ("Hand Weeding & Hoeing", "Physically pulling or cutting weeds using hand tools like a hoe or khurpi around crop plants."),
            ("Tractor Cultivators", "Using tractor-mounted rotary cultivators to uproot weeds quickly between crop rows."),
            ("Mowing & Slashing", "Cutting tall weeds before they flower, especially along field borders and orchard rows."),
            ("Soil Solarization", "Covering moist soil with clear plastic sheets during hot sunny weeks to heat the soil and kill weed seeds.")
        ]),
        ("7. Mulching Methods (Organic & Plastic)", [
            ("What is Mulching?", "Covering the bare soil around plants to block sunlight so weed seeds cannot sprout, while conserving moisture."),
            ("Organic Mulch", "Spreading dry straw, leaves, or crop waste. Naturally decomposes over time, enriching the soil with organic matter."),
            ("Plastic Mulch", "Laying black or silver polyethylene film over raised beds with planting holes. Widely used for vegetables and fruits."),
            ("Main Benefits", "Blocks 90%+ weed emergence, saves up to 50% irrigation water, and protects soil temperature.")
        ]),
        ("8. Biological Methods", [
            ("Beneficial Insects", "Releasing insects that feed specifically on target weeds without harming crops."),
            ("Natural Bio-Agents", "Using natural fungal or microbial agents that weaken specific weeds."),
            ("Grazing Animals & Ducks", "Using ducks in rice fields or grazing animals in orchards to eat weed sprouts."),
            ("Main Benefit", "100% eco-friendly, zero chemical residues, and completely safe for soil and pollinators.")
        ]),
        ("9. Advantages & Limitations", [
            ("Advantages", "Higher crop yields, efficient fertilizer and water use, cleaner harvest produce, fewer pests and diseases."),
            ("Limitations", "Hand weeding requires heavy labor; machinery and plastic mulch require upfront cost; wet weather delays weeding; plastic mulch requires proper collection.")
        ]),
        ("10. Conclusion: 5 Key Golden Rules", [
            ("Rule 1", "Start clean with certified seed and clean equipment."),
            ("Rule 2", "Act early in the first 30 days of crop growth."),
            ("Rule 3", "Protect the soil with organic or plastic mulch."),
            ("Rule 4", "Combine cultural, mechanical, and biological methods."),
            ("Rule 5", "Stop weeds before they flower and produce seeds.")
        ])
    ]

    for heading, items in sections_content:
        h = doc.add_heading(heading, level=1)
        h.style.font.color.rgb = DocxRGBColor(21, 128, 61)
        
        for item_title, item_desc in items:
            p = doc.add_paragraph()
            r_it = p.add_run(f"• {item_title}: ")
            r_it.bold = True
            r_it.font.color.rgb = DocxRGBColor(30, 41, 59)
            
            r_desc = p.add_run(item_desc)
            r_desc.font.color.rgb = DocxRGBColor(71, 85, 105)

    doc.add_paragraph("―" * 45)
    p_end = doc.add_paragraph()
    r_end = p_end.add_run("🌾 THANK YOU! Open for Discussion & Questions.")
    r_end.bold = True
    r_end.font.size = DocxPt(12)
    r_end.font.color.rgb = DocxRGBColor(21, 128, 61)

    doc.save("Weed_Management_Comprehensive_Guide.docx")
    print("Successfully generated Weed_Management_Comprehensive_Guide.docx")


def create_markdown_guide():
    md_content = """# Weed Management in Agriculture

> **Document Type:** Simple Presenter Handbook & Reference Guide  
> **Topic:** Meaning, Importance, Principles, Preventive, Cultural, Mechanical & Biological Methods, Mulching, Advantages, Limitations & Conclusion  

---

## 1. Meaning of Weed & Weed Management

### What is a Weed?
A weed is **any plant growing where it is not wanted**. For example, wild plants that grow in crop fields and compete with crops for food, water, space, and sunlight.

### What is Weed Management?
**Weed Management** is the practice of controlling and suppressing unwanted weeds so crops can grow healthy, receive full soil nutrients, and produce good harvest yields.

---

## 2. Importance of Weed Management

- **Protects Crop Yield:** Weeds rob crops of light, space, and nutrients. Controlling weeds prevents severe harvest losses.
- **Saves Water & Fertilizer:** Weeds absorb water and fertilizer rapidly. Managing weeds ensures these expensive inputs nourish your crops.
- **Prevents Pests & Crop Diseases:** Weeds serve as hiding and breeding spots for harmful insects, viruses, and plant pathogens.
- **Improves Produce Quality & Income:** Clean fields yield higher quality produce that sells at top market prices.

---

## 3. Principles of Weed Management

1. **Prevention First:** Stopping weed seeds from entering the field is much easier and cheaper than removing grown weeds.
2. **Critical Early Window:** Crops are most sensitive during the first 20 to 35 days. Keeping fields clean early allows crops to form a canopy that naturally shades out weeds.
3. **Combine Multiple Methods:** Use cultural, mechanical, and biological methods together instead of relying on only one tool.
4. **Stop Weed Seed Production:** Never let weeds flower and produce seeds. Removing weeds early reduces weed problems in future years.

---

## 4. Preventive Methods

- **Clean Certified Seeds:** Always sow certified crop seeds free from weed seeds.
- **Clean Farm Machinery:** Wash tractors and harvesting machinery before moving to clean fields.
- **Well-Rotted Compost:** Only use fully decomposed manure where weed seeds have been eliminated by natural heat.
- **Clean Field Borders:** Keep irrigation canals and field borders clear so seeds do not blow or wash into fields.

---

## 5. Cultural Methods

- **Stale Seedbed:** Irrigate the soil before sowing, let weeds sprout, kill them with light raking, and then sow your crop into a clean bed.
- **Crop Rotation:** Rotate different crops each season to disrupt the life cycles of weeds.
- **Proper Plant Spacing:** Plant crops at the recommended spacing so the crop canopy closes quickly and shades out weeds.
- **Cover Crops:** Grow fast-growing cover crops between rows to blanket the soil and naturally smother weeds.

---

## 6. Mechanical Methods

- **Hand Weeding & Hoeing:** Physically pulling or cutting weeds using simple hand tools like a hoe or *khurpi* around crop plants.
- **Tractor Cultivators:** Using tractor-mounted rotary cultivators to uproot weeds quickly between crop rows.
- **Mowing & Slashing:** Cutting tall weeds before they flower, especially along field borders and orchard rows.
- **Soil Solarization:** Covering moist soil with clear plastic sheets during hot sunny weeks to heat the soil and kill weed seeds.

---

## 7. Mulching Methods (Organic & Plastic Mulch)

- **What is Mulching?:** Covering the bare soil around plants to block sunlight so weed seeds cannot sprout, while conserving soil moisture.
- **Organic Mulch:** Spreading dry straw, leaves, or crop waste (naturally rots into organic matter that feeds the soil).
- **Plastic Mulch:** Laying black or silver polyethylene film over raised beds with planting holes (widely used for vegetables and fruits).
- **Key Benefits:** Over 90% weed control, saves up to 50% irrigation water, and protects soil temperature.

---

## 8. Biological Methods

- **Beneficial Insects:** Releasing specific insects that feed exclusively on harmful weeds without eating or damaging farm crops.
- **Natural Bio-Agents:** Using natural fungal or microbial agents that target and weaken specific weed species.
- **Grazing Animals & Ducks:** Using ducks in rice fields or grazing animals in orchards to eat weed sprouts naturally.
- **Main Benefit:** 100% eco-friendly, zero chemical residues, and completely safe for soil, pollinators, and people.

---

## 9. Advantages & Practical Limitations

| Aspect | Advantages | Practical Limitations |
| :--- | :--- | :--- |
| **Weed Management** | • Higher crop yields<br>• Saves water & fertilizers<br>• Clean produce quality<br>• Fewer pests & diseases | • Hand weeding requires heavy labor<br>• Machinery and plastic mulch require upfront cost<br>• Wet weather can delay physical weeding<br>• Plastic mulch requires disposal |

---

## 10. Conclusion: 5 Key Golden Rules

1. **Start Clean:** Always use certified clean seeds and sanitized equipment.
2. **Act Early:** Keep crops weed-free during the critical first 30 days of growth.
3. **Protect the Soil:** Use organic or plastic mulch to conserve water and block weeds.
4. **Combine Methods:** Integrate cultural, mechanical, and biological practices.
5. **Prevent Seed Setting:** Remove weeds before they flower to stop weed buildup in future years.

---

## 🌾 Presentation Closing

```
╔═══════════════════════════════════════════════════════════════════╗
║                                                                   ║
║                          🌾 THANK YOU! 🌾                         ║
║                                                                   ║
║            Thank you for your valuable time and attention!        ║
║                                                                   ║
║           Healthy Crops  •  Thriving Farms  •  Higher Yields      ║
║                                                                   ║
║             💬 Open for Questions, Comments & Discussion          ║
║                                                                   ║
╚═══════════════════════════════════════════════════════════════════╝
```
"""
    with open("Weed_Management_Guide.md", "w", encoding="utf-8") as f:
        f.write(md_content)
    print("Successfully generated Weed_Management_Guide.md")


def create_html_presentation():
    html_content = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=5.0, user-scalable=yes" />
  <meta name="color-scheme" content="only dark" />
  <meta name="theme-color" content="#0a130e" />
  <title>Weed Management in Agriculture — Slide Deck</title>
  <meta name="description" content="Meaning, Principles, Methods, Mulching, Advantages, Limitations & Conclusion." />
  
  <!-- Google Fonts -->
  <link rel="preconnect" href="https://fonts.googleapis.com" />
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin />
  <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@500;600;700;800&family=Plus+Jakarta+Sans:wght@400;500;600;700&display=swap" rel="stylesheet" />
  
  <style>
    /* Strict Mobile Dark Inversion Prevention */
    :root {
      forced-color-adjust: none !important;
      -webkit-print-color-adjust: exact !important;
      color-scheme: only dark !important;
      
      --bg-dark: #0a130e;
      --card-bg: #14241a;
      --card-border: #22c55e;
      --box-bg: #0f1d15;
      
      --accent-green: #22c55e;
      --accent-lime: #86efac;
      --accent-gold: #f59e0b;
      --accent-cyan: #06b6d4;
      --accent-red: #ef4444;
      
      --text-white: #ffffff;
      --text-muted: #e2e8f0;
      --text-dim: #94a3b8;
    }

    * {
      box-sizing: border-box;
      margin: 0;
      padding: 0;
      -webkit-tap-highlight-color: transparent;
      -webkit-text-size-adjust: 100% !important;
    }

    html, body {
      background-color: #060b08 !important;
      color: #ffffff !important;
      font-family: 'Plus Jakarta Sans', -apple-system, BlinkMacSystemFont, Roboto, sans-serif;
      min-height: 100vh;
      overflow-x: hidden;
      display: flex;
      flex-direction: column;
    }

    /* Ambient Animated Particle Canvas */
    #particles-canvas {
      position: fixed;
      top: 0;
      left: 0;
      width: 100vw;
      height: 100vh;
      pointer-events: none;
      z-index: 0;
      opacity: 0.6;
    }

    /* Top Progress Bar */
    .progress-bar-container {
      position: fixed;
      top: 0;
      left: 0;
      width: 100%;
      height: 4px;
      background: rgba(255, 255, 255, 0.1);
      z-index: 300;
    }

    .progress-bar {
      height: 100%;
      width: 8.33%;
      background: linear-gradient(90deg, #f59e0b, #22c55e);
      transition: width 0.4s cubic-bezier(0.16, 1, 0.3, 1);
      box-shadow: 0 0 10px #22c55e;
    }

    /* App Header */
    header {
      position: sticky;
      top: 0;
      z-index: 200;
      backdrop-filter: blur(20px);
      background: rgba(10, 19, 14, 0.95) !important;
      border-bottom: 1.5px solid rgba(34, 197, 94, 0.3);
      padding: 12px 20px;
      display: flex;
      align-items: center;
      justify-content: space-between;
      gap: 12px;
    }

    .brand {
      display: flex;
      align-items: center;
      gap: 10px;
    }

    .brand-logo {
      width: 38px;
      height: 38px;
      border-radius: 10px;
      background: linear-gradient(135deg, #15803d, #22c55e);
      display: flex;
      align-items: center;
      justify-content: center;
      font-size: 20px;
      box-shadow: 0 4px 12px rgba(34, 197, 94, 0.4);
      animation: logoPulse 3s ease-in-out infinite;
    }

    .brand-text h1 {
      font-family: 'Outfit', sans-serif;
      font-size: 1.15rem;
      font-weight: 700;
      color: #ffffff !important;
      line-height: 1.2;
    }

    .brand-text p {
      font-size: 0.75rem;
      color: #94a3b8 !important;
    }

    .header-btns {
      display: flex;
      align-items: center;
      gap: 8px;
    }

    .btn {
      padding: 8px 14px;
      border-radius: 8px;
      font-size: 0.85rem;
      font-weight: 600;
      cursor: pointer;
      display: inline-flex;
      align-items: center;
      gap: 6px;
      text-decoration: none;
      transition: all 0.2s ease;
      border: 1px solid transparent;
      user-select: none;
    }

    .btn-green {
      background: #22c55e !important;
      color: #05140a !important;
      font-weight: 700;
      box-shadow: 0 4px 12px rgba(34, 197, 94, 0.35);
    }
    .btn-green:hover {
      background: #16a34a !important;
      transform: translateY(-2px);
    }

    .btn-ghost {
      background: rgba(255, 255, 255, 0.08) !important;
      border-color: rgba(34, 197, 94, 0.35) !important;
      color: #ffffff !important;
    }
    .btn-ghost:hover {
      background: rgba(255, 255, 255, 0.16) !important;
    }

    /* Main Presentation Frame (16:9 Canvas Viewport) */
    main {
      flex: 1;
      width: 100%;
      max-width: 1220px;
      margin: 0 auto;
      padding: 16px 16px 24px;
      display: flex;
      flex-direction: column;
      gap: 14px;
      position: relative;
      z-index: 10;
    }

    /* 16:9 Slide Canvas (Maintains identical layout on mobile & desktop) */
    .deck-viewport {
      width: 100%;
      background: #14241a !important;
      border: 1.5px solid rgba(34, 197, 94, 0.4);
      border-radius: 16px;
      box-shadow: 0 20px 45px rgba(0, 0, 0, 0.6), 0 0 25px rgba(34, 197, 94, 0.12);
      min-height: 520px;
      padding: 24px 28px;
      display: flex;
      flex-direction: column;
      position: relative;
      overflow: hidden;
      transition: all 0.3s ease;
    }

    .slide {
      display: none;
      width: 100%;
      height: 100%;
    }

    .slide.active {
      display: flex;
      flex-direction: column;
      animation: slideEntrance 0.45s cubic-bezier(0.16, 1, 0.3, 1) forwards;
    }

    /* Slide Header */
    .slide-head {
      display: flex;
      justify-content: space-between;
      align-items: center;
      border-bottom: 1.5px solid rgba(34, 197, 94, 0.25);
      padding-bottom: 12px;
      margin-bottom: 18px;
      gap: 12px;
    }

    .slide-tag {
      font-size: 0.78rem;
      font-weight: 700;
      letter-spacing: 1.2px;
      color: #22c55e !important;
      text-transform: uppercase;
      margin-bottom: 2px;
    }

    .slide-title {
      font-family: 'Outfit', sans-serif;
      font-size: 1.65rem;
      font-weight: 700;
      color: #ffffff !important;
      line-height: 1.25;
    }

    .slide-badge {
      background: rgba(245, 158, 11, 0.15) !important;
      color: #f59e0b !important;
      border: 1px solid rgba(245, 158, 11, 0.5);
      padding: 5px 12px;
      border-radius: 20px;
      font-weight: 700;
      font-size: 0.82rem;
      white-space: nowrap;
    }

    /* Slide Body Layouts */
    .slide-content-2col {
      display: grid;
      grid-template-columns: 1.05fr 0.95fr;
      gap: 20px;
      align-items: center;
      flex: 1;
    }

    .slide-content-4col {
      display: grid;
      grid-template-columns: repeat(4, 1fr);
      gap: 14px;
      flex: 1;
    }

    .slide-content-2x2 {
      display: grid;
      grid-template-columns: 1fr 1fr;
      gap: 14px;
      flex: 1;
    }

    /* Photo Frame with Ambient Glow */
    .photo-frame {
      border-radius: 12px;
      overflow: hidden;
      border: 1.5px solid rgba(34, 197, 94, 0.35);
      height: 320px;
      background: #000;
      box-shadow: 0 10px 24px rgba(0, 0, 0, 0.5);
      position: relative;
    }

    .photo-frame img {
      width: 100%;
      height: 100%;
      object-fit: cover;
      transition: transform 0.5s ease;
    }

    .photo-frame:hover img {
      transform: scale(1.04);
    }

    /* Information Card Blocks */
    .info-card {
      background: #0f1d15 !important;
      border: 1.5px solid rgba(34, 197, 94, 0.3);
      border-radius: 12px;
      padding: 14px 16px;
      margin-bottom: 10px;
      box-shadow: 0 4px 12px rgba(0, 0, 0, 0.2);
      transition: all 0.25s ease;
    }

    .info-card:hover {
      transform: translateY(-2px);
      border-color: #22c55e !important;
      box-shadow: 0 6px 16px rgba(34, 197, 94, 0.18);
    }

    .info-title {
      font-size: 1.02rem;
      font-weight: 700;
      color: #22c55e !important;
      margin-bottom: 4px;
      display: flex;
      align-items: center;
      gap: 6px;
    }

    .info-desc {
      font-size: 0.92rem;
      color: #e2e8f0 !important;
      line-height: 1.5;
    }

    .info-desc strong {
      color: #ffffff !important;
    }

    /* Column Card (Slide 5) */
    .col-card-box {
      background: #0f1d15 !important;
      border: 1.5px solid rgba(34, 197, 94, 0.3);
      border-radius: 12px;
      padding: 16px 12px;
      display: flex;
      flex-direction: column;
      gap: 6px;
      transition: transform 0.25s ease;
    }

    .col-card-box:hover {
      transform: translateY(-3px);
      border-color: #f59e0b !important;
    }

    .col-icon {
      font-size: 1.7rem;
    }

    .col-title {
      font-size: 1rem;
      font-weight: 700;
      color: #f59e0b !important;
    }

    .col-desc {
      font-size: 0.88rem;
      color: #e2e8f0 !important;
      line-height: 1.4;
    }

    /* Presenter Notes Bar */
    .presenter-notes {
      background: #0b150f !important;
      border: 1px solid rgba(245, 158, 11, 0.4);
      border-radius: 10px;
      padding: 12px 16px;
      margin-top: 14px;
      animation: fadeIn 0.3s ease;
    }

    .notes-tag {
      font-size: 0.78rem;
      font-weight: 700;
      color: #f59e0b !important;
      text-transform: uppercase;
      margin-bottom: 3px;
    }

    .notes-text {
      font-size: 0.86rem;
      color: #cbd5e1 !important;
      line-height: 1.45;
    }

    /* Bottom Control Bar */
    .deck-controls {
      display: flex;
      justify-content: space-between;
      align-items: center;
      gap: 10px;
      flex-wrap: wrap;
    }

    .dots-list {
      display: flex;
      gap: 5px;
      flex-wrap: wrap;
      justify-content: center;
    }

    .dot-btn {
      width: 10px;
      height: 10px;
      border-radius: 50%;
      background: rgba(255, 255, 255, 0.2);
      cursor: pointer;
      transition: all 0.25s ease;
      border: none;
    }

    .dot-btn.active {
      background: #22c55e !important;
      width: 24px;
      border-radius: 10px;
      box-shadow: 0 0 10px #22c55e;
    }

    .control-actions {
      display: flex;
      align-items: center;
      gap: 8px;
    }

    /* Keyframe Animations */
    @keyframes slideEntrance {
      0% {
        opacity: 0;
        transform: translateY(14px) scale(0.985);
      }
      100% {
        opacity: 1;
        transform: translateY(0) scale(1);
      }
    }

    @keyframes fadeIn {
      from { opacity: 0; }
      to { opacity: 1; }
    }

    @keyframes logoPulse {
      0%, 100% { transform: translateY(0) scale(1); }
      50% { transform: translateY(-3px) scale(1.04); }
    }

    /* Confetti Canvas for Thank You Slide */
    #confetti-canvas {
      position: absolute;
      top: 0;
      left: 0;
      width: 100%;
      height: 100%;
      pointer-events: none;
      z-index: 50;
    }

    /* Compact Mobile Display without Breaking Slide Composition */
    @media (max-width: 820px) {
      header {
        padding: 10px 14px;
      }
      .brand-text p {
        display: none;
      }
      main {
        padding: 10px 8px 20px;
      }
      .deck-viewport {
        padding: 16px 14px;
        min-height: auto;
      }
      .slide-title {
        font-size: 1.25rem;
      }
      .slide-content-2col {
        grid-template-columns: 1fr;
        gap: 14px;
      }
      .slide-content-4col {
        grid-template-columns: 1fr 1fr;
        gap: 10px;
      }
      .slide-content-2x2 {
        grid-template-columns: 1fr;
        gap: 10px;
      }
      .photo-frame {
        height: 180px;
      }
      .info-card {
        padding: 10px 12px;
        margin-bottom: 8px;
      }
      .info-title {
        font-size: 0.95rem;
      }
      .info-desc {
        font-size: 0.85rem;
      }
      .deck-controls {
        flex-direction: column;
        gap: 10px;
      }
      .control-actions {
        width: 100%;
        justify-content: space-between;
      }
      .btn {
        flex: 1;
        justify-content: center;
      }
    }
  </style>
</head>
<body>

  <!-- Animated Ambient Particles -->
  <canvas id="particles-canvas"></canvas>

  <!-- Top Progress Bar -->
  <div class="progress-bar-container">
    <div class="progress-bar" id="progress-bar"></div>
  </div>

  <header>
    <div class="brand">
      <div class="brand-logo">🌱</div>
      <div class="brand-text">
        <h1>Weed Management</h1>
        <p>Simple Agronomy Guide</p>
      </div>
    </div>
    <div class="header-btns">
      <button class="btn btn-ghost" id="autoplay-btn" onclick="toggleAutoPlay()" title="Auto Play Slideshow">▶ Auto</button>
      <a href="Weed_Management_Mastery.pptx" download class="btn btn-green">⬇ PPTX</a>
      <a href="Weed_Management_Comprehensive_Guide.docx" download class="btn btn-ghost">📄 Word</a>
    </div>
  </header>

  <main>
    <div class="deck-viewport" id="slide-canvas">
      <canvas id="confetti-canvas"></canvas>
      
      <!-- Slide 1: Title -->
      <div class="slide active" id="slide-1">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Introduction</div>
            <h2 class="slide-title">Weed Management in Agriculture</h2>
          </div>
          <div class="slide-badge">01 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🌾 Agriculture & Crop Science</div>
              <p class="info-desc">A simple, complete guide to weed control techniques, principles, and methods.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #f59e0b;">📋 Key Topics Covered:</div>
              <p class="info-desc">
                • Meaning of Weed & Weed Management<br>
                • Importance & Core Principles<br>
                • Preventive, Cultural, Mechanical & Biological Methods<br>
                • Mulching Methods (Organic & Plastic)<br>
                • Advantages & Practical Limitations
              </p>
            </div>
          </div>
          <div class="photo-frame">
            <img src="assets/integrated_weed_management.jpg" alt="Weed Management" />
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Welcome everyone. Today we are presenting a complete guide to Weed Management in Agriculture. We will explain definitions, why controlling weeds is essential, core principles, all key methods, and advantages and limitations.</div>
        </div>
      </div>

      <!-- Slide 2: Meaning -->
      <div class="slide" id="slide-2">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Definition</div>
            <h2 class="slide-title">Meaning: What is a Weed & Weed Management?</h2>
          </div>
          <div class="slide-badge">02 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/weed_meaning_intro.jpg" alt="Weed in crop field" />
          </div>
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🌿 What is a Weed?</div>
              <p class="info-desc">A weed is <strong>any plant growing where it is not wanted</strong> (e.g. wild plants competing with farm crops for food, water, and sunlight).</p>
            </div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">🎯 What is Weed Management?</div>
              <p class="info-desc">The practice of controlling unwanted weeds so crops can grow healthy, receive full soil nutrients, and produce good harvest yields.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">⚠️ How Weeds Harm Crops</div>
              <p class="info-desc">Weeds steal sunlight, water, space, and fertilizers from the main crop, reducing overall harvest quality and quantity.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">What is a weed? Simply a plant out of place. Weed management is the organized way to suppress these unwanted plants so our crops thrive.</div>
        </div>
      </div>

      <!-- Slide 3: Importance -->
      <div class="slide" id="slide-3">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Importance</div>
            <h2 class="slide-title">Why Weed Management is Important</h2>
          </div>
          <div class="slide-badge">03 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/weed_crop_competition.jpg" alt="Crop Competition" />
          </div>
          <div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">📈 Protects Crop Yield & Harvest</div>
              <p class="info-desc">Weeds steal sunlight, space, and nutrients. Controlling weeds prevents large yield losses.</p>
            </div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">💧 Saves Water & Expensive Fertilizers</div>
              <p class="info-desc">Weeds absorb water and fertilizers rapidly. Management ensures inputs feed the crops, not weeds.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">🛡️ Prevents Pests & Crop Diseases</div>
              <p class="info-desc">Weeds act as breeding and hiding grounds for harmful insects, viruses, and fungi.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #22c55e;">💰 Higher Produce Quality & Profits</div>
              <p class="info-desc">Clean fields yield weed-free grain and vegetables that command premium market prices.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Why manage weeds? It protects harvest yield, ensures fertilizers nourish crops, stops pest infestation, and boosts farm profit.</div>
        </div>
      </div>

      <!-- Slide 4: Principles -->
      <div class="slide" id="slide-4">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Core Rules</div>
            <h2 class="slide-title">Core Principles of Weed Management</h2>
          </div>
          <div class="slide-badge">04 / 12</div>
        </div>
        <div class="slide-content-2x2">
          <div class="info-card" style="border-color: #22c55e;">
            <div class="info-title">🛡️ 1. Prevention First</div>
            <p class="info-desc">Stop weed seeds before they enter the farm. Keeping seeds out is much cheaper than removing grown weeds.</p>
          </div>
          <div class="info-card" style="border-color: #f59e0b;">
            <div class="info-title" style="color: #f59e0b;">⏱️ 2. Critical Early Window</div>
            <p class="info-desc">Keep crops weed-free during the first 20 to 35 days. After that, the crop canopy naturally shades out late weeds.</p>
          </div>
          <div class="info-card" style="border-color: #06b6d4;">
            <div class="info-title" style="color: #06b6d4;">🔄 3. Combine Multiple Methods</div>
            <p class="info-desc">Use cultural, mechanical, and biological tools together instead of relying on only one single method.</p>
          </div>
          <div class="info-card">
            <div class="info-title" style="color: #22c55e;">📉 4. Stop Weed Seed Production</div>
            <p class="info-desc">Never let weeds flower and produce seeds. Removing weeds early reduces weed pressure in future seasons.</p>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Remember the four core rules: Prevent seeds from entering, protect the first 30 days, integrate multiple tools, and never let weeds produce seeds.</div>
        </div>
      </div>

      <!-- Slide 5: Preventive Methods -->
      <div class="slide" id="slide-5">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Control Methods</div>
            <h2 class="slide-title">Preventive Methods: Stopping Weeds at the Gate</h2>
          </div>
          <div class="slide-badge">05 / 12</div>
        </div>
        <div class="slide-content-4col">
          <div class="col-card-box" style="border-color: #22c55e;">
            <div class="col-icon">🌾</div>
            <div class="col-title" style="color: #22c55e;">Clean Certified Seeds</div>
            <div class="col-desc">Always sow certified crop seeds free from weed seed contamination.</div>
          </div>
          <div class="col-card-box" style="border-color: #f59e0b;">
            <div class="col-icon">🚜</div>
            <div class="col-title">Machinery Hygiene</div>
            <div class="col-desc">Wash tractors and harvesting equipment before moving to clean fields.</div>
          </div>
          <div class="col-card-box" style="border-color: #06b6d4;">
            <div class="col-icon">🧪</div>
            <div class="col-title" style="color: #06b6d4;">Well-Rotted Compost</div>
            <div class="col-desc">Only apply decomposed manure where weed seeds have been destroyed by natural composting heat.</div>
          </div>
          <div class="col-card-box">
            <div class="col-icon">🌊</div>
            <div class="col-title" style="color: #22c55e;">Clean Irrigation</div>
            <div class="col-desc">Keep canals and field borders clean so weed seeds do not wash or blow into fields.</div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Preventive methods are the most economical step: clean seeds, clean machinery, fully decomposed manure, and clean irrigation canals.</div>
        </div>
      </div>

      <!-- Slide 6: Cultural Methods -->
      <div class="slide" id="slide-6">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Control Methods</div>
            <h2 class="slide-title">Cultural Methods: Agronomic Farming Practices</h2>
          </div>
          <div class="slide-badge">06 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/cover_crop_mulching.jpg" alt="Cover crops" />
          </div>
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🌱 Stale Seedbed Preparation</div>
              <p class="info-desc">Water the field early, let weeds sprout, kill them with light raking, then sow your crop in a clean bed.</p>
            </div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">🔄 Crop Rotation</div>
              <p class="info-desc">Alternate different crops each season to disrupt the life cycles of specific weeds.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">📐 Proper Spacing & Density</div>
              <p class="info-desc">Plant crops at recommended spacing so the crop canopy closes fast and shades out weeds.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #22c55e;">🌿 Cover Crops</div>
              <p class="info-desc">Grow fast-growing cover crops between rows to blanket the soil and smother weeds naturally.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Cultural methods use good farming practices: stale seedbeds, rotating crops, proper spacing for quick shading, and cover crops.</div>
        </div>
      </div>

      <!-- Slide 7: Mechanical Methods -->
      <div class="slide" id="slide-7">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Control Methods</div>
            <h2 class="slide-title">Mechanical Methods: Physical Tools & Machinery</h2>
          </div>
          <div class="slide-badge">07 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/mechanical_weeding_tractor.jpg" alt="Tractor Weeding" />
          </div>
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🛠️ Hand Weeding & Hoeing</div>
              <p class="info-desc">Physically pulling or cutting weeds using hand tools like a hoe or khurpi around crop plants.</p>
            </div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">🚜 Tractor Inter-Row Cultivators</div>
              <p class="info-desc">Using tractor-mounted rotary weeders to uproot weeds quickly between crop rows.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">✂️ Mowing & Slashing</div>
              <p class="info-desc">Cutting tall weeds before flowering along borders and orchard alleys.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #22c55e;">☀️ Soil Solarization</div>
              <p class="info-desc">Covering moist soil with clear plastic sheets in hot sun to kill topsoil weed seeds.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Mechanical weeding uses physical force: hand hoes for delicate work, tractor cultivators for large acreage, and mowing before weeds flower.</div>
        </div>
      </div>

      <!-- Slide 8: Mulching Methods -->
      <div class="slide" id="slide-8">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Soil Cover</div>
            <h2 class="slide-title">Mulching Methods: Organic & Plastic Covers</h2>
          </div>
          <div class="slide-badge">08 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/organic_plastic_mulching.jpg" alt="Mulching" />
          </div>
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🛡️ What is Mulching?</div>
              <p class="info-desc">Covering bare soil to block sunlight so weed seeds cannot sprout, while conserving moisture.</p>
            </div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">🌾 Organic Mulch (Straw & Leaves)</div>
              <p class="info-desc">Using straw or crop waste that naturally breaks down into soil organic matter.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">🖤 Plastic Mulch (Black Film)</div>
              <p class="info-desc">Laying black plastic sheets over raised beds; widely used in vegetable and fruit farming.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #22c55e;">⭐ Triple Benefit</div>
              <p class="info-desc">Blocks 90%+ weeds, saves 50% water, and protects soil temperature.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Mulch works by blocking light so weed seeds cannot germinate. Organic straw enriches soil, while plastic mulch saves huge amounts of water.</div>
        </div>
      </div>

      <!-- Slide 9: Biological Methods -->
      <div class="slide" id="slide-9">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Eco-Friendly</div>
            <h2 class="slide-title">Biological Methods: Natural Weed Control</h2>
          </div>
          <div class="slide-badge">09 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div class="photo-frame">
            <img src="assets/biological_control_insect.jpg" alt="Bio Control" />
          </div>
          <div>
            <div class="info-card" style="border-color: #22c55e;">
              <div class="info-title">🐞 Beneficial Insects</div>
              <p class="info-desc">Releasing insects that feed exclusively on specific problem weeds without touching crops.</p>
            </div>
            <div class="info-card" style="border-color: #f59e0b;">
              <div class="info-title" style="color: #f59e0b;">🍄 Natural Bio-Agents & Fungi</div>
              <p class="info-desc">Using targeted natural fungi or bio-agents that weaken specific weed plants.</p>
            </div>
            <div class="info-card" style="border-color: #06b6d4;">
              <div class="info-title" style="color: #06b6d4;">🦆 Grazing Animals & Ducks</div>
              <p class="info-desc">Using ducks in rice fields or grazing animals in orchards to eat weeds naturally.</p>
            </div>
            <div class="info-card">
              <div class="info-title" style="color: #22c55e;">🌿 100% Eco-Friendly</div>
              <p class="info-desc">Zero chemical residues, safe for soil, pollinators, and farm workers.</p>
            </div>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Biological control uses nature's balance: specific insects, natural fungi, or animals like ducks to suppress weeds without chemicals.</div>
        </div>
      </div>

      <!-- Slide 10: Advantages & Limitations -->
      <div class="slide" id="slide-10">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Comparison</div>
            <h2 class="slide-title">Advantages & Practical Limitations</h2>
          </div>
          <div class="slide-badge">10 / 12</div>
        </div>
        <div class="slide-content-2x2">
          <div class="info-card" style="border-color: #22c55e;">
            <div class="info-title">✅ Major Advantages</div>
            <p class="info-desc" style="line-height: 1.8;">
              ✔ <strong>Higher Crop Yields:</strong> Full sunlight and nutrients.<br>
              ✔ <strong>Saves Water & Fertilizer:</strong> Inputs nourish the crop.<br>
              ✔ <strong>Better Produce Quality:</strong> Clean harvest commands top price.<br>
              ✔ <strong>Fewer Pests & Diseases:</strong> Breaks pest hiding spots.<br>
              ✔ <strong>Long-term Soil Health:</strong> Reduces future weed seeds.
            </p>
          </div>
          <div class="info-card" style="border-color: #ef4444;">
            <div class="info-title" style="color: #ef4444;">⚠️ Practical Limitations</div>
            <p class="info-desc" style="line-height: 1.8;">
              ✖ <strong>High Labor Demand:</strong> Hand weeding is labor-intensive.<br>
              ✖ <strong>Equipment Costs:</strong> Tractors and mulch need upfront money.<br>
              ✖ <strong>Weather Dependence:</strong> Heavy rain delays mechanical weeding.<br>
              ✖ <strong>Plastic Waste:</strong> Plastic mulch needs proper disposal after harvest.
            </p>
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">A balanced view: Weed management greatly increases yield and quality, but requires planning for labor and equipment costs.</div>
        </div>
      </div>

      <!-- Slide 11: Conclusion -->
      <div class="slide" id="slide-11">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Summary</div>
            <h2 class="slide-title">Conclusion: 5 Key Golden Rules</h2>
          </div>
          <div class="slide-badge">11 / 12</div>
        </div>
        <div class="info-card" style="border-color: #22c55e;">
          <div class="info-title">🎯 The 5 Golden Takeaways:</div>
          <p class="info-desc" style="line-height: 2; font-size: 1.02rem;">
            <strong>1. Start Clean:</strong> Always use certified clean seeds and sanitized equipment.<br>
            <strong>2. Act Early:</strong> Keep crops weed-free during the critical first 30 days of growth.<br>
            <strong>3. Protect the Soil:</strong> Use organic or plastic mulch to conserve water and block weeds.<br>
            <strong>4. Combine Methods:</strong> Integrate cultural, mechanical, and biological practices.<br>
            <strong>5. Prevent Seed Setting:</strong> Remove weeds before they flower to stop future weed buildup.
          </p>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">To conclude, remember: Start clean, act early, use mulch, combine multiple tactics, and stop weeds before they seed.</div>
        </div>
      </div>

      <!-- Slide 12: Thank You -->
      <div class="slide" id="slide-12">
        <div class="slide-head">
          <div>
            <div class="slide-tag">• Closing</div>
            <h2 class="slide-title">Thank You!</h2>
          </div>
          <div class="slide-badge">12 / 12</div>
        </div>
        <div class="slide-content-2col">
          <div>
            <div class="info-card" style="border-color: #f59e0b; padding: 24px 18px;">
              <h2 style="color: #f59e0b; font-size: 2.2rem; margin-bottom: 10px;">🌾 THANK YOU!</h2>
              <p style="font-size: 1.1rem; margin-bottom: 12px;"><strong>Thank you for your valuable time and attention.</strong></p>
              <p style="font-size: 0.98rem; color: #22c55e; font-weight: 700; margin-bottom: 16px;">Healthy Crops • Thriving Farms • Higher Yields</p>
              <p style="font-size: 0.92rem; color: #94a3b8;">💬 Questions, comments, and discussion are welcome!</p>
            </div>
          </div>
          <div class="photo-frame">
            <img src="assets/conclusion_thankyou_farm.jpg" alt="Farm Sunset" />
          </div>
        </div>
        <div class="presenter-notes">
          <div class="notes-tag">🎤 Presenter Notes</div>
          <div class="notes-text">Thank you everyone for listening. We now open the floor for questions and discussion.</div>
        </div>
      </div>

    </div>

    <!-- Navigation Bar -->
    <div class="deck-controls">
      <button class="btn btn-ghost" id="prev-btn" onclick="prevSlide()">◀ Previous</button>
      <div class="dots-list" id="dots-container"></div>
      <div class="control-actions">
        <button class="btn btn-green" id="next-btn" onclick="nextSlide()">Next Slide ▶</button>
      </div>
    </div>
  </main>

  <script>
    let currentSlide = 1;
    const totalSlides = 12;
    let autoPlayTimer = null;

    // Subtle audio synthesis for slide change
    const audioCtx = (window.AudioContext || window.webkitAudioContext) ? new (window.AudioContext || window.webkitAudioContext)() : null;
    function playSlideSound() {
      if (!audioCtx) return;
      try {
        if (audioCtx.state === 'suspended') audioCtx.resume();
        const osc = audioCtx.createOscillator();
        const gain = audioCtx.createGain();
        osc.type = 'sine';
        osc.frequency.setValueAtTime(520, audioCtx.currentTime);
        osc.frequency.exponentialRampToValueAtTime(880, audioCtx.currentTime + 0.08);
        gain.gain.setValueAtTime(0.04, audioCtx.currentTime);
        gain.gain.exponentialRampToValueAtTime(0.001, audioCtx.currentTime + 0.08);
        osc.connect(gain);
        gain.connect(audioCtx.destination);
        osc.start();
        osc.stop(audioCtx.currentTime + 0.08);
      } catch (e) {}
    }

    function initDots() {
      const container = document.getElementById('dots-container');
      container.innerHTML = '';
      for (let i = 1; i <= totalSlides; i++) {
        const dot = document.createElement('button');
        dot.className = `dot-btn ${i === 1 ? 'active' : ''}`;
        dot.title = `Slide ${i}`;
        dot.onclick = () => showSlide(i);
        container.appendChild(dot);
      }
    }

    function showSlide(n) {
      if (n < 1) n = 1;
      if (n > totalSlides) n = totalSlides;
      currentSlide = n;

      for (let i = 1; i <= totalSlides; i++) {
        const slide = document.getElementById(`slide-${i}`);
        if (slide) slide.classList.remove('active');
      }

      const activeSlide = document.getElementById(`slide-${currentSlide}`);
      if (activeSlide) activeSlide.classList.add('active');

      const dots = document.querySelectorAll('.dot-btn');
      dots.forEach((dot, idx) => {
        dot.classList.toggle('active', idx + 1 === currentSlide);
      });

      const progress = (currentSlide / totalSlides) * 100;
      document.getElementById('progress-bar').style.width = `${progress}%`;

      document.getElementById('prev-btn').style.visibility = currentSlide === 1 ? 'hidden' : 'visible';
      document.getElementById('next-btn').innerText = currentSlide === totalSlides ? 'Finished 🎉' : 'Next Slide ▶';

      playSlideSound();

      if (currentSlide === totalSlides) {
        launchConfetti();
      }
    }

    function prevSlide() {
      showSlide(currentSlide - 1);
    }

    function nextSlide() {
      if (currentSlide < totalSlides) {
        showSlide(currentSlide + 1);
      } else {
        showSlide(1);
      }
    }

    // Touch Swipe Gesture Handling
    let touchStartX = 0;
    let touchEndX = 0;
    const slideCanvas = document.getElementById('slide-canvas');

    slideCanvas.addEventListener('touchstart', (e) => {
      touchStartX = e.changedTouches[0].screenX;
    }, { passive: true });

    slideCanvas.addEventListener('touchend', (e) => {
      touchEndX = e.changedTouches[0].screenX;
      const diff = touchEndX - touchStartX;
      if (Math.abs(diff) > 40) {
        if (diff < 0) nextSlide();
        else prevSlide();
      }
    }, { passive: true });

    // Keyboard navigation
    document.addEventListener('keydown', (e) => {
      if (e.key === 'ArrowRight' || e.key === 'Space') nextSlide();
      else if (e.key === 'ArrowLeft') prevSlide();
    });

    // Auto Play Slideshow
    function toggleAutoPlay() {
      const btn = document.getElementById('autoplay-btn');
      if (autoPlayTimer) {
        clearInterval(autoPlayTimer);
        autoPlayTimer = null;
        btn.innerText = '▶ Auto';
        btn.classList.remove('btn-green');
        btn.classList.add('btn-ghost');
      } else {
        autoPlayTimer = setInterval(nextSlide, 5000);
        btn.innerText = '⏸ Pause';
        btn.classList.remove('btn-ghost');
        btn.classList.add('btn-green');
      }
    }

    // Ambient Particle Background Animation
    function initParticles() {
      const canvas = document.getElementById('particles-canvas');
      const ctx = canvas.getContext('2d');
      let w = canvas.width = window.innerWidth;
      let h = canvas.height = window.innerHeight;

      window.addEventListener('resize', () => {
        w = canvas.width = window.innerWidth;
        h = canvas.height = window.innerHeight;
      });

      const particles = [];
      for (let i = 0; i < 35; i++) {
        particles.push({
          x: Math.random() * w,
          y: Math.random() * h,
          radius: Math.random() * 2 + 1,
          vx: (Math.random() - 0.5) * 0.4,
          vy: -Math.random() * 0.5 - 0.2,
          color: Math.random() > 0.4 ? 'rgba(34, 197, 94, ' : 'rgba(245, 158, 11, ',
          alpha: Math.random() * 0.5 + 0.2
        });
      }

      function render() {
        ctx.clearRect(0, 0, w, h);
        particles.forEach(p => {
          p.x += p.vx;
          p.y += p.vy;
          if (p.y < 0) p.y = h;
          if (p.x < 0) p.x = w;
          if (p.x > w) p.x = 0;

          ctx.beginPath();
          ctx.arc(p.x, p.y, p.radius, 0, Math.PI * 2);
          ctx.fillStyle = p.color + p.alpha + ')';
          ctx.shadowBlur = 8;
          ctx.shadowColor = '#22c55e';
          ctx.fill();
        });
        requestAnimationFrame(render);
      }
      render();
    }

    // Confetti on Slide 12
    function launchConfetti() {
      const c = document.getElementById('confetti-canvas');
      const ctx = c.getContext('2d');
      c.width = c.offsetWidth;
      c.height = c.offsetHeight;

      const pieces = [];
      const colors = ['#22c55e', '#f59e0b', '#06b6d4', '#ffffff', '#86efac'];
      for (let i = 0; i < 60; i++) {
        pieces.push({
          x: c.width / 2,
          y: c.height / 2,
          vx: (Math.random() - 0.5) * 12,
          vy: (Math.random() - 0.5) * 12 - 3,
          size: Math.random() * 6 + 4,
          color: colors[Math.floor(Math.random() * colors.length)],
          rotation: Math.random() * 360,
          vRot: (Math.random() - 0.5) * 10,
          opacity: 1
        });
      }

      let frame = 0;
      function step() {
        ctx.clearRect(0, 0, c.width, c.height);
        pieces.forEach(p => {
          p.x += p.vx;
          p.y += p.vy;
          p.vy += 0.25; // gravity
          p.rotation += p.vRot;
          p.opacity -= 0.012;

          if (p.opacity > 0) {
            ctx.save();
            ctx.translate(p.x, p.y);
            ctx.rotate((p.rotation * Math.PI) / 180);
            ctx.fillStyle = p.color;
            ctx.globalAlpha = Math.max(0, p.opacity);
            ctx.fillRect(-p.size / 2, -p.size / 2, p.size, p.size);
            ctx.restore();
          }
        });

        frame++;
        if (frame < 90) requestAnimationFrame(step);
        else ctx.clearRect(0, 0, c.width, c.height);
      }
      step();
    }

    initDots();
    initParticles();
    showSlide(1);
  </script>
</body>
</html>
"""
    with open("index.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    print("Successfully generated index.html (Fixed Phone & Laptop 1:1 Layout with Animations)")

if __name__ == "__main__":
    create_pptx()
    create_docx()
    create_markdown_guide()
    create_html_presentation()
    print("\nALL DELIVERABLES REBUILT IDENTICALLY FOR PHONE & LAPTOP WITH ATTRACTIVE ANIMATIONS!")
